"""
Build the Astronomer SE Lab demo presentation.
Run: python3 build_presentation.py
Output: Astronomer_SE_Lab_Demo.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colors ──────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0C, 0x15, 0x4A)   # deep Astronomer navy
PURPLE = RGBColor(0x7B, 0x5C, 0xEA)   # Astronomer purple accent
BLUE   = RGBColor(0x00, 0xA6, 0xE2)   # sky blue
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF4, 0xF4, 0xF8)   # near-white bg
GRAY   = RGBColor(0x6B, 0x7B, 0x8D)
GREEN  = RGBColor(0x22, 0xC5, 0x5E)
ORANGE = RGBColor(0xF5, 0xA6, 0x23)
RED    = RGBColor(0xEF, 0x44, 0x44)

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # truly blank layout


# ══════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill=NAVY, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_para(tf, text, size=16, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def slide_header(slide, title, subtitle=None, accent=PURPLE):
    """Standard header bar across the top."""
    add_rect(slide, 0, 0, 13.33, 1.3, fill=NAVY)
    add_rect(slide, 0, 1.25, 13.33, 0.08, fill=accent)
    add_text(slide, title, 0.4, 0.12, 11, 0.7,
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.78, 11, 0.45,
                 size=15, color=BLUE)


def screenshot_box(slide, x, y, w, h, label="[ Screenshot ]"):
    """Dashed placeholder box for screenshots."""
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xE8, 0xEA, 0xF2)
    box.line.color.rgb = PURPLE
    box.line.width = Pt(1.5)
    # label
    txb = slide.shapes.add_textbox(
        Inches(x), Inches(y + h/2 - 0.25), Inches(w), Inches(0.5))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(13)
    run.font.color.rgb = GRAY
    run.font.italic = True
    run.font.name = "Calibri"


def pill(slide, x, y, w, h, text, fill=PURPLE, text_color=WHITE, size=12):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Calibri"


def arrow_right(slide, x, y, length=0.45):
    """Simple right-pointing arrow line."""
    connector = slide.shapes.add_connector(
        1,
        Inches(x), Inches(y),
        Inches(x + length), Inches(y)
    )
    connector.line.color.rgb = PURPLE
    connector.line.width = Pt(2)


def flow_box(slide, x, y, w, h, title, body_lines, fill=NAVY,
             title_color=WHITE, body_color=WHITE, title_size=13, body_size=11):
    add_rect(slide, x, y, w, h, fill=fill)
    # title
    txb = slide.shapes.add_textbox(
        Inches(x + 0.1), Inches(y + 0.08), Inches(w - 0.2), Inches(0.35))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = title_color
    run.font.name = "Calibri"
    # body
    if body_lines:
        txb2 = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.43),
            Inches(w - 0.2), Inches(h - 0.5))
        txb2.word_wrap = True
        tf2 = txb2.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(body_lines):
            p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = line
            run2.font.size = Pt(body_size)
            run2.font.color.rgb = body_color
            run2.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 0, 4.6, 13.33, 0.1, fill=PURPLE)
add_rect(sl, 0, 4.7, 13.33, 2.8, fill=RGBColor(0x0A, 0x10, 0x38))

add_text(sl, "ASTRONOMER", 0.5, 0.5, 12, 0.8,
         size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(sl, "SE Lab — Solution Walkthrough", 0.5, 1.4, 12.33, 1.5,
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Building, Deploying & Monitoring Airflow Pipelines on Astro",
         0.5, 2.95, 12.33, 0.7,
         size=22, bold=False, color=BLUE, align=PP_ALIGN.CENTER)
add_text(sl, "Michael Gnesin  ·  Sales Engineering Candidate  ·  May 2026",
         0.5, 5.15, 12.33, 0.5,
         size=15, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT)
slide_header(sl, "Agenda", "15-Minute Walkthrough")

items = [
    ("1", "Part 1 — Creating DAGs",
     "Serial ETL · Parallel DQ Checks · Crazy Complex · External Connections"),
    ("2", "Part 2 — Git Integration & CI/CD",
     "GitHub repo · Astronomer Git deploy · push-to-deploy flow"),
    ("3", "Part 3 — Running & Monitoring",
     "Hosted deployment · Worker Queues · Alerting"),
    ("4", "Recap & Q&A",
     "Architecture summary · Key takeaways"),
]

for i, (num, title, sub) in enumerate(items):
    y = 1.6 + i * 1.35
    add_rect(sl, 0.4, y, 0.55, 0.9, fill=PURPLE)
    add_text(sl, num, 0.4, y + 0.15, 0.55, 0.6,
             size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 1.1, y, 11.8, 0.9, fill=WHITE)
    add_text(sl, title, 1.25, y + 0.05, 6, 0.4,
             size=18, bold=True, color=NAVY)
    add_text(sl, sub, 1.25, y + 0.45, 11, 0.35,
             size=13, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture Overview
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT)
slide_header(sl, "Solution Architecture", "End-to-end data platform on Astronomer")

# Source box
flow_box(sl, 0.35, 1.55, 2.0, 2.0, "Source",
         ["GCP Cloud SQL", "PostgreSQL", "(app_data.customers)"],
         fill=RGBColor(0x1D, 0x4E, 0x89))
# Arrow
arrow_right(sl, 2.38, 2.55, 0.55)
# DAGs box
flow_box(sl, 2.95, 1.55, 3.5, 2.0, "Airflow on Astro",
         ["dag_serial_etl", "dag_2_parallel_dq_checks", "dag_3_crazy_complex"],
         fill=PURPLE)
# Arrow
arrow_right(sl, 6.48, 2.55, 0.55)
# Snowflake box
flow_box(sl, 7.05, 1.55, 2.0, 2.0, "Destination",
         ["Snowflake", "AUDIT.PIPELINE_LOG", "(DEMO_DB)"],
         fill=RGBColor(0x29, 0xB5, 0xE8))
# Arrow down from Airflow
add_rect(sl, 4.65, 3.6, 0.08, 0.55, fill=PURPLE)
# Monitoring box
flow_box(sl, 3.2, 4.2, 2.9, 1.5, "Observability",
         ["Astro Alerts", "Task-level failure notifications", "Email / PagerDuty"],
         fill=RGBColor(0x0F, 0x3D, 0x2E))

# Right side — key features
add_rect(sl, 9.5, 1.55, 3.5, 5.4, fill=NAVY)
add_text(sl, "Key Capabilities", 9.6, 1.65, 3.3, 0.4,
         size=14, bold=True, color=BLUE)
features = [
    "TaskFlow API (Python decorators)",
    "Shared Snowflake Connection",
    "Git-based CI/CD deploy",
    "Worker Queue routing",
    "Scheduled execution (daily 6AM)",
    "Automated DQ assertions",
    "DAG-level & task-level alerts",
    "XCom data passing between tasks",
]
txb = sl.shapes.add_textbox(Inches(9.6), Inches(2.15), Inches(3.2), Inches(4.5))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
for j, f in enumerate(features):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = "✓  " + f
    run.font.size = Pt(12)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Part 1 section divider
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 0, 3.3, 13.33, 0.12, fill=PURPLE)
add_text(sl, "PART 1", 0.5, 1.8, 12.33, 0.9,
         size=18, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
add_text(sl, "Creating DAGs", 0.5, 2.5, 12.33, 1.2,
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Serial ETL  ·  Parallel DQ  ·  Crazy Complex  ·  External Connections",
         0.5, 4.0, 12.33, 0.6,
         size=18, color=BLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DAG 1: Serial ETL
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT)
slide_header(sl, "DAG 1: Serial ETL Pipeline",
             "dag_serial_etl  ·  GCP PostgreSQL → Snowflake  ·  TaskFlow API")

# Task flow boxes
tasks = [
    ("extract_from\n_postgres", "SELECT COUNT(*)\nFROM app_data\n.customers", RGBColor(0x1D,0x4E,0x89)),
    ("validate\n_data",         "Abort if\ncustomer_count\n== 0",              RGBColor(0x6C,0x35,0xB5)),
    ("transform\n_data",        "Enrich payload\nwith timestamp\n& pipeline name", PURPLE),
    ("load_to\n_snowflake",     "INSERT into\nAUDIT.PIPELINE\n_LOG",           RGBColor(0x29,0xB5,0xE8)),
]

box_w = 2.2
gap   = 0.55
start_x = 0.35
y_box = 1.65

for i, (name, desc, col) in enumerate(tasks):
    x = start_x + i * (box_w + gap)
    flow_box(sl, x, y_box, box_w, 2.0, name, desc.split("\n"),
             fill=col, title_size=13, body_size=11)
    if i < len(tasks) - 1:
        ax = x + box_w + 0.05
        arrow_right(sl, ax, y_box + 1.0, gap - 0.12)

# Serial dependency callout
add_rect(sl, 0.35, 3.85, 11.83, 0.55, fill=RGBColor(0xE8,0xEA,0xF2))
add_text(sl, "Each task passes its output to the next via XCom — strict serial execution enforced by TaskFlow API return values.",
         0.55, 3.88, 11.5, 0.5, size=12, color=NAVY)

# Tags / schedule badge
pill(sl, 0.35, 4.6, 1.3, 0.38, "tags: serial, se-lab", fill=NAVY, size=11)
pill(sl, 1.75, 4.6, 1.5, 0.38, "schedule: manual", fill=GRAY, size=11)
pill(sl, 3.35, 4.6, 2.0, 0.38, "connection: snowflake", fill=RGBColor(0x29,0xB5,0xE8), size=11)
pill(sl, 5.45, 4.6, 2.2, 0.38, "connection: gcp_postgres", fill=RGBColor(0x1D,0x4E,0x89), size=11)

# Code snippet
add_rect(sl, 0.35, 5.15, 7.8, 2.0, fill=RGBColor(0x1E,0x1E,0x2E))
code = (
    "@task()\n"
    "def extract_from_postgres():\n"
    "    hook = PostgresHook(postgres_conn_id='gcp_postgres')\n"
    "    result = hook.get_records('SELECT COUNT(*) FROM app_data.customers;')\n"
    "    return {'customer_count': result[0][0]}\n\n"
    "@task()\n"
    "def load_to_snowflake(payload: dict):\n"
    "    hook = SnowflakeHook(snowflake_conn_id='snowflake')\n"
    "    hook.run(f\"INSERT INTO AUDIT.PIPELINE_LOG ...\")"
)
add_text(sl, code, 0.5, 5.2, 7.6, 1.9,
         size=10, color=RGBColor(0xA8,0xFF,0xC2), bold=False)

# Screenshot placeholder
screenshot_box(sl, 8.3, 5.1, 4.7, 2.1, "[ Screenshot: Graph View ]")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DAG 2: Parallel DQ Checks
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT)
slide_header(sl, "DAG 2: Parallel Data Quality Checks",
             "dag_2_parallel_dq_checks  ·  Snowflake  ·  Scheduled daily at 6 AM UTC")

# 4 parallel task boxes
dq_tasks = [
    ("check_null\n_customer_emails",  "NULL EMAIL\ncheck",            "Fails if any\ncustomer row\nlacks email",     RGBColor(0x6C,0x35,0xB5)),
    ("check_orphaned\n_orders",       "REFERENTIAL\nINTEGRITY",       "Detects orders\nwith no matching\ncustomer",  RGBColor(0x1D,0x4E,0x89)),
    ("check_negative\n_order_amounts","AMOUNT\nRANGE",                "Flags rows\nwhere amount\n< 0",               RGBColor(0x29,0x78,0xC2)),
    ("check_duplicate\n_customers",   "UNIQUENESS\ncheck",            "GROUP BY email\nHAVING count > 1",            PURPLE),
]

# Trigger node
add_rect(sl, 0.35, 2.5, 1.5, 0.6, fill=NAVY)
add_text(sl, "Schedule\nTrigger", 0.35, 2.5, 1.5, 0.6,
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box_w = 2.4
gap   = 0.25
start_x = 2.2
y_box = 1.6

for i, (name, badge, desc, col) in enumerate(dq_tasks):
    x = start_x + i * (box_w + gap)
    # line from trigger
    arrow_right(sl, 1.87, 2.8, x - 1.89)
    flow_box(sl, x, y_box, box_w, 2.2, name,
             [badge, "", desc], fill=col, title_size=12, body_size=10)

# All 4 run in parallel callout
add_rect(sl, 2.2, 3.95, 11.0, 0.5, fill=RGBColor(0xE8,0xEA,0xF2))
add_text(sl, "All 4 tasks execute simultaneously — no dependencies. Each uses SnowflakeHook and raises AssertionError on failure.",
         2.4, 3.98, 10.6, 0.45, size=12, color=NAVY)

# Tags / schedule
pill(sl, 0.35, 4.65, 2.0, 0.38, "schedule: 0 6 * * *", fill=NAVY, size=11)
pill(sl, 2.45, 4.65, 2.3, 0.38, "tags: parallel, data-quality", fill=PURPLE, size=11)
pill(sl, 4.85, 4.65, 2.2, 0.38, "connection: snowflake", fill=RGBColor(0x29,0xB5,0xE8), size=11)

# Screenshot
screenshot_box(sl, 0.35, 5.2, 6.1, 2.0, "[ Screenshot: Parallel Graph View ]")
screenshot_box(sl, 6.65, 5.2, 6.3, 2.0, "[ Screenshot: Successful DAG Run ]")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DAG 3: Crazy Complex
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT)
slide_header(sl, "DAG 3: The Crazy Complex DAG",
             "dag_3_crazy_complex  ·  7 layers  ·  20+ tasks  ·  Maximum graph complexity")

# Layers description — left column
add_rect(sl, 0.35, 1.55, 5.5, 5.4, fill=NAVY)
layers = [
    ("Layer 0", "start",                          "Entry point"),
    ("Layer 1", "extract_customers/orders/products", "Fan-out: 3 extractors"),
    ("Layer 2", "validate_* (×6)",                "2 validators per source"),
    ("Layer 3", "join_cust_orders, join_orders_products, join_all",
                                                  "Cross-stream joins → complex edges"),
    ("Layer 4", "transform_revenue/cohorts/inventory/clv",
                                                  "Parallel transforms"),
    ("Layer 5", "quality_gate",                   "Convergence — all 4 transforms required"),
    ("Layer 6", "load_snowflake/gcs/reporting_mart", "Fan-out: 3 parallel loads"),
    ("Layer 7", "notify_* · update_data_catalog", "Fan-in → end"),
]
txb = sl.shapes.add_textbox(Inches(0.45), Inches(1.65), Inches(5.3), Inches(5.1))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
for j, (lyr, tasks_l, desc) in enumerate(layers):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.space_before = Pt(5)
    run = p.add_run()
    run.text = f"{lyr}:  "
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = BLUE
    run.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(11)
    run2.font.color.rgb = WHITE
    run2.font.name = "Calibri"

# Why it's complex
add_rect(sl, 0.35, 6.15, 5.5, 0.7, fill=PURPLE)
add_text(sl, "Cross-layer data sharing creates an irregular, highly-connected graph — not a simple tree.",
         0.45, 6.17, 5.3, 0.65, size=11, color=WHITE)

# Screenshot placeholder — large
screenshot_box(sl, 6.1, 1.55, 6.85, 5.4, "[ Screenshot: Graph View — full DAG ]")
screenshot_box(sl, 6.1, 7.05, 6.85, 0.3, "")  # not needed, skip

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — External Connections
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT)
slide_header(sl, "Connecting to External Systems",
             "Shared Connections configured once in Astro — reused across all DAGs")

# Two connection cards
for ci, (icon, name, conn_id, details, col) in enumerate([
    ("❄️", "Snowflake", "snowflake",
     ["Used by: dag_serial_etl, dag_2_parallel_dq_checks, dag_3_crazy_complex",
      "Writes audit logs to AUDIT.PIPELINE_LOG",
      "Runs DQ assertions on DEMO_DB tables",
      "Connection type: Snowflake  (account, warehouse, role, database)"],
     RGBColor(0x29,0xB5,0xE8)),
    ("🐘", "GCP PostgreSQL", "gcp_postgres",
     ["Used by: dag_serial_etl",
      "Source: app_data.customers table",
      "Query: SELECT COUNT(*) FROM app_data.customers",
      "Connection type: Postgres  (host, port, schema, login, password)"],
     RGBColor(0x1D,0x4E,0x89)),
]):
    x = 0.35 + ci * 6.5
    add_rect(sl, x, 1.55, 6.1, 3.8, fill=col)
    add_text(sl, name, x + 0.2, 1.65, 5.7, 0.55,
             size=22, bold=True, color=WHITE)
    add_text(sl, f"conn_id: {conn_id}", x + 0.2, 2.2, 5.7, 0.4,
             size=13, color=LIGHT)
    txb = sl.shapes.add_textbox(
        Inches(x + 0.2), Inches(2.65), Inches(5.6), Inches(2.4))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for k, d in enumerate(details):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = "•  " + d
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"

# Screenshot
screenshot_box(sl, 0.35, 5.55, 12.6, 1.65,
               "[ Screenshot: Astro Connections UI — showing both connections ]")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Part 2 section divider
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 0, 3.3, 13.33, 0.12, fill=BLUE)
add_text(sl, "PART 2", 0.5, 1.8, 12.33, 0.9,
         size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(sl, "Git Integration & CI/CD", 0.5, 2.5, 12.33, 1.2,
         size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "GitHub  ·  Astronomer Git Deploy  ·  Push-to-Deploy",
         0.5, 4.0, 12.33, 0.6,
         size=18, color=BLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Git Integration
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT)
slide_header(sl, "Git Integration & CI/CD Pipeline",
             "github.com/mgnesin/se-lab  →  Astronomer Hosted Deployment")

# Flow diagram
steps = [
    ("1", "Local Dev",      "astro dev start\nEdit DAGs\nRun tests",       NAVY),
    ("2", "git push\nmain", "Triggers GitHub\nintegration",                PURPLE),
    ("3", "Astronomer\nGit Deploy", "Detects push\nBuilds image\nDeploys", RGBColor(0x0F,0x6B,0xC4)),
    ("4", "Astro\nDeployment", "New DAGs live\nin minutes",                RGBColor(0x29,0xB5,0xE8)),
]

bw = 2.3
bh = 2.0
sy = 2.0
sx = 0.4

for i, (num, title, desc, col) in enumerate(steps):
    x = sx + i * (bw + 0.55)
    add_rect(sl, x, sy, bw, bh, fill=col)
    add_text(sl, num, x + 0.1, sy + 0.1, 0.5, 0.45,
             size=22, bold=True, color=BLUE)
    add_text(sl, title, x + 0.1, sy + 0.5, bw - 0.2, 0.55,
             size=16, bold=True, color=WHITE)
    add_text(sl, desc, x + 0.1, sy + 1.05, bw - 0.2, 0.85,
             size=12, color=LIGHT)
    if i < len(steps) - 1:
        ax = x + bw + 0.05
        arrow_right(sl, ax, sy + 1.0, 0.42)

# Key points
add_rect(sl, 0.4, 4.3, 12.6, 0.55, fill=NAVY)
add_text(sl, "  Workflow: git commit → git push → Astronomer detects change → auto build & deploy → DAGs live in Airflow UI",
         0.5, 4.33, 12.3, 0.48, size=13, color=WHITE)

# Commit history note
add_rect(sl, 0.4, 5.05, 5.8, 1.85, fill=NAVY)
add_text(sl, "Recent Commits (main branch)", 0.55, 5.12, 5.5, 0.4,
         size=13, bold=True, color=BLUE)
commits = [
    "8090ac9  Test CI/CD deploy",
    "dabbc85  Test CI/CD deploy",
    "a8375a1  Test CI/CD deploy",
    "Initial  Astro SE Lab project",
]
txb = sl.shapes.add_textbox(Inches(0.55), Inches(5.58), Inches(5.5), Inches(1.2))
tf = txb.text_frame
for k, c in enumerate(commits):
    p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = c
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0xA8,0xFF,0xC2)
    run.font.name = "Courier New"

# Screenshot
screenshot_box(sl, 6.45, 5.0, 6.5, 1.95,
               "[ Screenshot: Astro Git Integration UI ]")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Part 3 section divider
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 0, 3.3, 13.33, 0.12, fill=GREEN)
add_text(sl, "PART 3", 0.5, 1.8, 12.33, 0.9,
         size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(sl, "Running & Monitoring", 0.5, 2.5, 12.33, 1.2,
         size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Hosted Deployment  ·  Worker Queues  ·  Alerts",
         0.5, 4.0, 12.33, 0.6,
         size=18, color=GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Running on Astro Hosted
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT)
slide_header(sl, "Running DAGs on Astro Hosted Deployment",
             "All 3 DAGs deployed and executed on Astronomer Cloud")

# 3 status cards
for ci, (dag_id, schedule, last_run, status, col) in enumerate([
    ("dag_serial_etl",           "Manual trigger", "May 12, 2026", "Success", GREEN),
    ("dag_2_parallel_dq_checks", "Daily @ 6AM UTC","May 13, 2026", "Success", GREEN),
    ("dag_3_crazy_complex",      "Manual trigger", "May 12, 2026", "Success", GREEN),
]):
    x = 0.35 + ci * 4.35
    add_rect(sl, x, 1.55, 4.0, 2.4, fill=NAVY)
    add_text(sl, dag_id, x + 0.15, 1.65, 3.7, 0.45,
             size=13, bold=True, color=BLUE)
    add_text(sl, f"Schedule:  {schedule}", x + 0.15, 2.15, 3.7, 0.35,
             size=12, color=WHITE)
    add_text(sl, f"Last Run:  {last_run}", x + 0.15, 2.52, 3.7, 0.35,
             size=12, color=WHITE)
    add_rect(sl, x + 0.15, 3.05, 1.4, 0.38, fill=col)
    add_text(sl, f"● {status}", x + 0.15, 3.05, 1.4, 0.38,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Screenshots
screenshot_box(sl, 0.35, 4.15, 6.1, 3.1,
               "[ Screenshot: Astro Deployment — DAG list view ]")
screenshot_box(sl, 6.65, 4.15, 6.3, 3.1,
               "[ Screenshot: Successful DAG run — task grid ]")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Worker Queues
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT)
slide_header(sl, "Worker Queues",
             "Routing compute-intensive tasks to a dedicated worker pool")

# Explanation
add_rect(sl, 0.35, 1.55, 12.6, 1.0, fill=NAVY)
add_text(sl, "Worker queues allow specific tasks to be routed to dedicated worker nodes — "
             "useful for isolating resource-intensive workloads (heavy transforms, ML inference) "
             "from lightweight tasks (API calls, DQ checks).",
         0.55, 1.6, 12.2, 0.9, size=13, color=WHITE)

# Two-queue diagram
# Default queue
add_rect(sl, 0.35, 2.75, 5.8, 3.5, fill=NAVY)
add_text(sl, "Default Worker Queue", 0.5, 2.82, 5.5, 0.45,
         size=16, bold=True, color=WHITE)
add_text(sl, "(Standard compute — all tasks unless overridden)",
         0.5, 3.25, 5.5, 0.35, size=11, color=GRAY)
default_tasks = [
    "extract_from_postgres",
    "validate_data",
    "check_null_customer_emails",
    "check_orphaned_orders",
    "check_negative_order_amounts",
    "check_duplicate_customers",
]
txb = sl.shapes.add_textbox(Inches(0.55), Inches(3.65), Inches(5.4), Inches(2.3))
tf = txb.text_frame
for k, t in enumerate(default_tasks):
    p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
    p.space_before = Pt(5)
    run = p.add_run()
    run.text = "▸  " + t
    run.font.size = Pt(12)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

# Heavy queue
add_rect(sl, 6.55, 2.75, 6.4, 3.5, fill=PURPLE)
add_text(sl, "heavy-transform Queue", 6.7, 2.82, 6.1, 0.45,
         size=16, bold=True, color=WHITE)
add_text(sl, "(High-memory workers — routed via queue='heavy-transform')",
         6.7, 3.25, 6.1, 0.35, size=11, color=LIGHT)
heavy_tasks = [
    "transform_data          queue='heavy-transform'",
    "load_to_snowflake       queue='heavy-transform'",
]
txb2 = sl.shapes.add_textbox(Inches(6.7), Inches(3.65), Inches(6.0), Inches(1.2))
tf2 = txb2.text_frame
for k, t in enumerate(heavy_tasks):
    p = tf2.paragraphs[0] if k == 0 else tf2.add_paragraph()
    p.space_before = Pt(8)
    run = p.add_run()
    run.text = "▸  " + t
    run.font.size = Pt(12)
    run.font.color.rgb = WHITE
    run.font.name = "Courier New"

add_rect(sl, 6.55, 4.85, 6.4, 1.2, fill=RGBColor(0x3D,0x1A,0x78))
add_text(sl, "Why?  The transform and load tasks are the most compute-intensive steps in the ETL "
             "pipeline. Isolating them prevents resource contention with lightweight DQ checks.",
         6.7, 4.9, 6.1, 1.1, size=12, color=WHITE)

# Screenshot
screenshot_box(sl, 0.35, 6.45, 12.6, 0.85,
               "[ Screenshot: Astro Worker Queue configuration + task detail showing queue assignment ]")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Alerts
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT)
slide_header(sl, "Alerting & Monitoring",
             "Proactive notifications on DAG execution conditions")

# Alert type cards
alerts = [
    ("DAG Failure Alert",
     "dag_serial_etl",
     "Triggers when any task in the serial ETL\npipeline fails (e.g. empty customer table,\nSnowflake write error).",
     "Notification: Email to data-team@example.com",
     RED),
    ("Task Duration Alert",
     "dag_2_parallel_dq_checks",
     "Triggers when any DQ check task exceeds\n60 seconds — signals a slow Snowflake query\nor warehouse scaling event.",
     "Notification: Slack #data-alerts channel",
     ORANGE),
    ("DAG Success Alert",
     "dag_serial_etl",
     "Confirms the nightly ETL completed\nsuccessfully. Useful for SLA tracking\nand stakeholder reporting.",
     "Notification: Email to analytics@example.com",
     GREEN),
]

for ci, (title, dag, desc, notif, col) in enumerate(alerts):
    x = 0.35 + ci * 4.35
    add_rect(sl, x, 1.6, 4.0, 3.7, fill=NAVY)
    add_rect(sl, x, 1.6, 4.0, 0.45, fill=col)
    add_text(sl, title, x + 0.12, 1.65, 3.8, 0.38,
             size=14, bold=True, color=WHITE)
    add_text(sl, f"DAG: {dag}", x + 0.12, 2.15, 3.8, 0.35,
             size=11, color=BLUE)
    add_text(sl, desc, x + 0.12, 2.55, 3.8, 1.2,
             size=12, color=WHITE)
    add_rect(sl, x + 0.12, 3.85, 3.76, 0.6, fill=RGBColor(0x1A,0x2A,0x4A))
    add_text(sl, notif, x + 0.22, 3.9, 3.56, 0.5,
             size=11, color=col)

# Why alerts matter
add_rect(sl, 0.35, 5.5, 12.6, 0.75, fill=PURPLE)
add_text(sl, "Alerts transform Airflow from a scheduler into an active monitoring system — "
             "teams are notified before failures compound into data quality incidents.",
         0.55, 5.55, 12.2, 0.65, size=13, color=WHITE)

# Screenshot
screenshot_box(sl, 0.35, 6.4, 12.6, 0.85,
               "[ Screenshot: Astro Alerts UI — showing configured alert rules ]")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Recap
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT)
slide_header(sl, "Recap & Key Takeaways", "What was built across Parts 1–3")

sections = [
    ("Part 1\nDAGs", [
        "Serial ETL: Postgres → validate → transform → Snowflake",
        "Parallel DQ: 4 Snowflake assertions run simultaneously",
        "Crazy Complex: 7-layer, 20+ task DAG with cross-joins",
        "Shared connections: Snowflake + GCP PostgreSQL",
    ], PURPLE),
    ("Part 2\nGit CI/CD", [
        "GitHub repo: github.com/mgnesin/se-lab",
        "Astronomer Git Integration configured",
        "Push-to-deploy: git push → Astro builds & deploys",
        "DAGs as code — full version history",
    ], RGBColor(0x0F,0x6B,0xC4)),
    ("Part 3\nMonitoring", [
        "All DAGs running on Astro Hosted Deployment",
        "Worker queues: heavy-transform pool for ETL tasks",
        "3 alert types: failure, duration, success",
        "Notifications to Email + Slack",
    ], GREEN),
]

for ci, (title, points, col) in enumerate(sections):
    x = 0.35 + ci * 4.35
    add_rect(sl, x, 1.55, 4.0, 5.6, fill=NAVY)
    add_rect(sl, x, 1.55, 4.0, 0.75, fill=col)
    add_text(sl, title, x + 0.15, 1.58, 3.7, 0.72,
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb = sl.shapes.add_textbox(
        Inches(x + 0.15), Inches(2.45), Inches(3.7), Inches(4.5))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for k, pt in enumerate(points):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = "✓  " + pt
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Q&A / Thank you
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 0, 5.3, 13.33, 0.1, fill=PURPLE)
add_text(sl, "ASTRONOMER", 0.5, 0.6, 12.33, 0.6,
         size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(sl, "Questions?", 0.5, 1.5, 12.33, 1.5,
         size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Michael Gnesin  ·  mgnesin@gmail.com",
         0.5, 3.2, 12.33, 0.6,
         size=20, color=GRAY, align=PP_ALIGN.CENTER)
add_text(sl, "github.com/mgnesin/se-lab",
         0.5, 3.85, 12.33, 0.5,
         size=16, color=BLUE, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/michaelgnesin/Projects/Astronomer/se-lab/Astronomer_SE_Lab_Demo.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
